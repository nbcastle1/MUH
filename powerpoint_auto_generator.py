#!/usr/bin/env python3
"""
Automated PowerPoint Generator for Motor Learning Stride Change Images
Creates a PowerPoint presentation with one image per slide, ordered by participant age.

Requirements:
- pip install python-pptx pandas

Usage:
1. Update the file paths in the main() function
2. Run the script: python powerpoint_generator.py
"""

import os
import re
import pandas as pd
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def extract_subject_id_from_filename(filename: str) -> Optional[str]:
    """
    Extract subject ID from filename.
    
    Examples:
    - "stride_change_MUH1396_fixed_grid.png" -> "MUH1396"
    - "stride_change_MUH1069_fixed_grid.png" -> "MUH1069"
    """
    # Match pattern: stride_change_[SUBJECT_ID]_fixed_grid.png
    match = re.search(r'stride_change_([A-Z]+\d+)_fixed_grid\.png', filename)
    if match:
        return match.group(1)
    return None

def load_subject_metadata(metadata_path: str) -> Dict[str, float]:
    """
    Load subject metadata and extract age information.
    Returns dictionary mapping subject_id -> age_in_years
    """
    try:
        # Load metadata CSV
        metadata_df = pd.read_csv(metadata_path)
        
        # Convert age_months to years if that column exists
        if 'age_months' in metadata_df.columns:
            metadata_df['age_years'] = metadata_df['age_months'] / 12
        elif 'age' in metadata_df.columns:
            metadata_df['age_years'] = metadata_df['age']
        else:
            print("⚠️ Warning: No age column found in metadata")
            return {}
        
        # Create mapping from ID to age
        age_mapping = {}
        for _, row in metadata_df.iterrows():
            subject_id = str(row['ID'])
            age = row.get('age_years', None)
            if pd.notna(age):
                age_mapping[subject_id] = float(age)
        
        print(f"✅ Loaded age data for {len(age_mapping)} subjects")
        return age_mapping
        
    except FileNotFoundError:
        print(f"❌ Metadata file not found: {metadata_path}")
        return {}
    except Exception as e:
        print(f"❌ Error loading metadata: {e}")
        return {}

def parse_image_paths_from_file(file_path: str) -> List[str]:
    """
    Parse image paths from your paste.txt file.
    """
    try:
        with open(file_path, 'r') as f:
            content = f.read()
        
        paths = []
        for line in content.strip().split('\n'):
            # Remove quotes and clean the path
            path = line.strip().strip('"')
            if path and path.endswith('.png'):
                paths.append(path)
        
        print(f"📁 Found {len(paths)} image paths in {file_path}")
        return paths
        
    except FileNotFoundError:
        print(f"❌ Image paths file not found: {file_path}")
        return []
    except Exception as e:
        print(f"❌ Error reading image paths: {e}")
        return []

def find_images_in_directory(directory: str, pattern: str = "*stride_change*fixed_grid.png") -> List[str]:
    """
    Alternative method: Find all stride change images in a directory.
    """
    try:
        directory_path = Path(directory)
        if not directory_path.exists():
            print(f"❌ Directory not found: {directory}")
            return []
        
        # Find all matching images
        image_paths = list(directory_path.rglob(pattern))
        image_paths_str = [str(path) for path in image_paths]
        
        print(f"📁 Found {len(image_paths_str)} images in {directory}")
        return image_paths_str
        
    except Exception as e:
        print(f"❌ Error searching directory: {e}")
        return []

def create_powerpoint_from_images(
    image_paths: List[str], 
    output_filename: str = "stride_change_presentation.pptx",
    metadata_path: Optional[str] = None,
    title: str = "Motor Learning Stride Change Analysis"
) -> str:
    """
    Create PowerPoint presentation from stride change images, ordered by age.
    
    Parameters:
    -----------
    image_paths : List[str]
        List of image file paths
    output_filename : str
        Name of output PowerPoint file
    metadata_path : str, optional
        Path to metadata CSV file containing age information
    title : str
        Title for the presentation
    
    Returns:
    --------
    str : Path to created PowerPoint file
    """
    
    print(f"🚀 Creating PowerPoint presentation: {output_filename}")
    
    # Create new presentation
    prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Load subject metadata for age information
    age_mapping = {}
    if metadata_path and os.path.exists(metadata_path):
        age_mapping = load_subject_metadata(metadata_path)
    
    # Process image paths and extract subject information
    image_data = []
    missing_images = []
    
    for img_path in image_paths:
        # Convert Windows path to current OS format and handle relative paths
        img_path = img_path.replace('\\', os.sep).replace('/', os.sep)
        
        # If path starts with C:\ but we're not on Windows, convert to relative path
        if img_path.startswith('C:') and os.name != 'nt':
            # Extract just the filename and look for it in current directory structure
            filename = os.path.basename(img_path)
            # Try to find the file in common locations
            possible_paths = [
                filename,
                f"figures/individual_plots/stride_change_after_success_vs_failure/{filename}",
                f"analysis/figures/individual_plots/stride_change_after_success_vs_failure/{filename}",
                f"motor_learning_output/figures/individual_plots/stride_change_after_success_vs_failure/{filename}"
            ]
            
            found = False
            for possible_path in possible_paths:
                if os.path.exists(possible_path):
                    img_path = possible_path
                    found = True
                    break
            
            if not found:
                print(f"⚠️ Could not find image: {filename}")
                continue
        
        # Extract subject ID from filename
        filename = os.path.basename(img_path)
        subject_id = extract_subject_id_from_filename(filename)
        
        if subject_id is None:
            print(f"⚠️ Could not extract subject ID from: {filename}")
            continue
        
        # Get age if available
        age = age_mapping.get(subject_id, None)
        
        # Check if image file exists
        if os.path.exists(img_path):
            image_data.append({
                'path': img_path,
                'subject_id': subject_id,
                'age': age,
                'filename': filename
            })
        else:
            missing_images.append({
                'path': img_path,
                'subject_id': subject_id,
                'age': age,
                'filename': filename
            })
    
    print(f"📊 Found {len(image_data)} valid images")
    if missing_images:
        print(f"⚠️ {len(missing_images)} images not found on disk")
        print("   First few missing:")
        for missing in missing_images[:5]:
            print(f"     {missing['filename']}")
    
    if not image_data:
        print("❌ No valid images found! Check your paths.")
        return ""
    
    # Sort by age (subjects without age data will be at the end)
    image_data.sort(key=lambda x: (x['age'] is None, x['age'] if x['age'] is not None else 999, x['subject_id']))
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Set title
    title_slide.shapes.title.text = title
    
    # Set subtitle with summary information
    if title_slide.placeholders[1]:
        subtitle_text = f"Individual Stride Change Distributions\n"
        subtitle_text += f"{len(image_data)} participants"
        if age_mapping:
            ages_with_data = [img['age'] for img in image_data if img['age'] is not None]
            if ages_with_data:
                subtitle_text += f"\nAge range: {min(ages_with_data):.1f} - {max(ages_with_data):.1f} years"
                subtitle_text += f"\nOrdered from youngest to oldest"
        
        title_slide.placeholders[1].text = subtitle_text
    
    # Add image slides
    blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
    
    print(f"📝 Adding {len(image_data)} image slides...")
    
    for i, img_info in enumerate(image_data, 1):
        # Create new slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.margin_left = Inches(0.1)
        title_frame.margin_right = Inches(0.1)
        title_frame.margin_top = Inches(0.1)
        title_frame.margin_bottom = Inches(0.1)
        
        # Set title text
        title_text = f"Subject {img_info['subject_id']}"
        if img_info['age'] is not None:
            title_text += f" (Age: {img_info['age']:.1f} years)"
        title_text += f" - Slide {i} of {len(image_data)}"
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title_text
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Format title text
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(47, 84, 150)  # Dark blue
        
        # Add image
        try:
            # Calculate image dimensions to fit slide while maintaining aspect ratio
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Leave space for title (1.2 inches from top) and small margin at bottom
            available_height = slide_height - Inches(1.2) - Inches(0.3)  # Title space + margin
            available_width = slide_width - Inches(1.0)  # Side margins
            
            # Add image with height constraint only (maintains aspect ratio)
            left = Inches(0.5)  # Starting position (will center later)
            top = Inches(1.2)   # Below title
            
            # Add the image with height only - PowerPoint will maintain aspect ratio
            pic = slide.shapes.add_picture(
                img_info['path'], 
                left, 
                top, 
                height=available_height
                # Note: No width specified - this maintains aspect ratio
            )
            
            # Center the image horizontally if it's narrower than available width
            actual_width = pic.width
            if actual_width < available_width:
                # Center horizontally in the available space
                horizontal_center = Inches(0.5) + (available_width - actual_width) / 2
                pic.left = int(horizontal_center)
            
            if i % 10 == 0:  # Progress update every 10 slides
                print(f"   Added slide {i}/{len(image_data)}")
                
        except Exception as e:
            print(f"❌ Error adding image {img_info['filename']}: {e}")
            # Add error text instead of image
            error_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(2))
            error_frame = error_box.text_frame
            error_paragraph = error_frame.paragraphs[0]
            error_paragraph.text = f"Error loading image:\n{img_info['filename']}"
            error_paragraph.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    try:
        prs.save(output_filename)
        print(f"✅ PowerPoint presentation saved: {output_filename}")
        print(f"📊 Total slides: {len(prs.slides)} (1 title + {len(image_data)} image slides)")
        
        # Print age summary
        if age_mapping:
            ages_with_data = [img['age'] for img in image_data if img['age'] is not None]
            ages_without_data = len([img for img in image_data if img['age'] is None])
            
            if ages_with_data:
                print(f"👥 Age range: {min(ages_with_data):.1f} - {max(ages_with_data):.1f} years")
                print(f"📈 Subjects with age data: {len(ages_with_data)}")
            if ages_without_data > 0:
                print(f"⚠️ Subjects without age data: {ages_without_data} (placed at end)")
        
        return output_filename
        
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return ""

def main():
    """
    Main function - Update these paths for your setup
    """
    
    print("🎯 Motor Learning Stride Change PowerPoint Generator")
    print("=" * 60)
    
    # =============================================================================
    # UPDATE THESE PATHS FOR YOUR SETUP
    # =============================================================================
    
    # Option 1: Use the paste.txt file with image paths
    image_paths_file = "paste.txt"  # Your file with the image paths
    
    # Option 2: Search a directory for images (alternative)
    # images_directory = "analysis/figures/individual_plots/stride_change_after_success_vs_failure/"
    
    # Metadata file with age information
    metadata_file = "muh_metadata.csv"
    
    # Output PowerPoint filename
    output_file = "stride_change_analysis_by_age.pptx"
    
    # Presentation title
    presentation_title = "Motor Learning: Stride Change After Success vs Failure"
    
    # =============================================================================
    
    # Method 1: Load paths from file
    if os.path.exists(image_paths_file):
        print(f"📁 Loading image paths from: {image_paths_file}")
        image_paths = parse_image_paths_from_file(image_paths_file)
    else:
        # Method 2: Search directory (fallback)
        print(f"📁 Image paths file not found, searching current directory...")
        image_paths = find_images_in_directory(".", "*stride_change*fixed_grid.png")
    
    if not image_paths:
        print("❌ No image paths found!")
        print("💡 Make sure either:")
        print("   1. paste.txt exists with your image paths, OR")
        print("   2. Image files are in the current directory or subdirectories")
        return
    
    # Create PowerPoint presentation
    result = create_powerpoint_from_images(
        image_paths=image_paths,
        output_filename=output_file,
        metadata_path=metadata_file,
        title=presentation_title
    )
    
    if result:
        print("\n" + "=" * 60)
        print("🎉 SUCCESS! PowerPoint presentation created!")
        print(f"📄 File: {result}")
        print("💡 Images are ordered from youngest to oldest participant")
        print("💡 Open the file in PowerPoint to view/edit")
    else:
        print("\n❌ Failed to create PowerPoint presentation")

if __name__ == "__main__":
    main()
